import os
import zipfile
import shutil
from pathlib import Path
from typing import List, Tuple

# ── Supported CODE extensions ──────────────────────────────────────────────────
SUPPORTED_CODE_EXTENSIONS = {
    '.py', '.js', '.jsx', '.ts', '.tsx',
    '.java', '.c', '.cpp', '.cc', '.cxx', '.h', '.hpp', '.cs',
    '.go', '.rs', '.rb', '.php',
    '.xml', '.yaml', '.yml',
    '.sh', '.bash',
    '.sql', '.html', '.css', '.scss',
    '.properties', '.gradle', '.pom'
}

# ── Supported DOCUMENT extensions (PDF, Office, etc.) ──────────────────────────
SUPPORTED_DOCUMENT_EXTENSIONS = {
    '.pdf',
    '.doc', '.docx',
    '.ppt', '.pptx',
    '.xls', '.xlsx', '.xlsm', '.csv',
    '.odt', '.ods', '.odp',
}

SUPPORTED_EXTENSIONS = SUPPORTED_CODE_EXTENSIONS | SUPPORTED_DOCUMENT_EXTENSIONS


# ── FIX 1: Noise files jo retrieval kharab karte hain — inhe SKIP karo ────────
# Yeh files har query pe retrieve hoti thi aur score giraati thi
SKIP_FILENAMES = {
    # Lock files — sirf dependency versions hain, code nahi
    'package-lock.json',
    'yarn.lock',
    'composer.lock',
    'poetry.lock',
    'Pipfile.lock',
    'Gemfile.lock',
    'pnpm-lock.yaml',
    'packages.lock.json',

    # Config/meta files — project ka logic nahi hota inme
    'package.json',        # sirf dependencies list
    '.gitignore',
    '.gitattributes',
    '.prettierrc',
    '.eslintrc',
    '.eslintignore',
    '.babelrc',
    '.editorconfig',
    '.browserslistrc',
    'tsconfig.json',
    'jsconfig.json',
    'jest.config.js',
    'webpack.config.js',
    'vite.config.js',
    'vite.config.ts',
    'tailwind.config.js',
    'postcss.config.js',
    'babel.config.js',
    'rollup.config.js',
    '.env.example',
    '.env.sample',
    'Makefile',
    'LICENSE',
    'CHANGELOG.md',
    'CONTRIBUTING.md',
    'CODE_OF_CONDUCT.md',

    # README — documentation hai, code nahi
    # Note: Agar tum chahte ho README embed ho toh yahan se hata do
    'README.md',
    'readme.md',
    'README.rst',
    'README.txt',
}

# ── FIX 1b: Noise file PATTERNS — extension ke basis pe skip karo ─────────────
SKIP_EXTENSIONS = {
    '.lock',       # sab lock files
    '.log',        # log files
    '.map',        # source map files (minified JS)
    '.min.js',     # minified JS — readable nahi
    '.chunk.js',   # webpack chunks
    '.snap',       # jest snapshots
    '.ico',        # icons
    '.png', '.jpg', '.jpeg', '.gif', '.svg', '.webp',  # images
    '.woff', '.woff2', '.ttf', '.eot',  # fonts
    '.zip', '.tar', '.gz',  # archives
    '.pyc',        # compiled python
}


# ──────────────────────────────────────────────────────────────────────────────
# ZIP extraction
# ──────────────────────────────────────────────────────────────────────────────
def extract_zip_file(zip_path: str, extract_to: str) -> str:
    try:
        if not zipfile.is_zipfile(zip_path):
            raise ValueError("Invalid ZIP file")

        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            for member in zip_ref.infolist():
                member_path = os.path.normpath(member.filename)
                if member_path.startswith('..') or os.path.isabs(member_path):
                    print(f"  [ZIP] Skipping unsafe path: {member.filename}", flush=True)
                    continue
                zip_ref.extract(member, extract_to)

        return extract_to
    except Exception as e:
        raise ValueError(f"Failed to extract ZIP: {str(e)}")


# ──────────────────────────────────────────────────────────────────────────────
# File discovery — FIX 1 applied here
# ──────────────────────────────────────────────────────────────────────────────
def get_supported_files(directory: str) -> List[Tuple[str, str]]:
    supported_files = []
    skipped_noise  = 0

    for root, dirs, files in os.walk(directory):
        # Skip irrelevant directories
        dirs[:] = [d for d in dirs if d not in {
            '.git', '.env', '__pycache__', 'node_modules',
            '.venv', 'venv', 'dist', 'build', '.next', '.nuxt',
            '.idea', '.vscode', 'target', 'coverage', '.pytest_cache',
            'vendor', 'bower_components', '.cache', 'tmp', 'temp',
            'logs', 'log', '.nyc_output', 'storybook-static'
        }]

        for file in files:
            file_lower = file.lower()
            file_ext   = Path(file).suffix.lower()

            # ── FIX 1: Skip noise filenames exactly ──────────────────────────
            if file in SKIP_FILENAMES or file_lower in SKIP_FILENAMES:
                skipped_noise += 1
                continue

            # ── FIX 1b: Skip noise extensions ────────────────────────────────
            if file_ext in SKIP_EXTENSIONS:
                skipped_noise += 1
                continue

            # ── Only keep supported code/document extensions ──────────────────
            if file_ext in SUPPORTED_EXTENSIONS:
                abs_path = os.path.join(root, file)
                rel_path = os.path.relpath(abs_path, directory)
                supported_files.append((abs_path, rel_path))

    if skipped_noise > 0:
        print(f"  [FILTER] Skipped {skipped_noise} noise files (lock/config/readme/images)", flush=True)

    return supported_files


# ──────────────────────────────────────────────────────────────────────────────
# Document readers — each returns plain text or ""
# ──────────────────────────────────────────────────────────────────────────────

def _read_pdf(file_path: str) -> str:
    try:
        from pdfminer.high_level import extract_text as pdfminer_extract
        text = pdfminer_extract(file_path)
        if text and text.strip():
            return text.strip()
    except ImportError:
        pass
    except Exception as e:
        print(f"  [PDF] pdfminer failed: {e}", flush=True)

    try:
        import PyPDF2
        text_parts = []
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                try:
                    t = page.extract_text()
                    if t:
                        text_parts.append(t)
                except Exception:
                    pass
        return '\n'.join(text_parts).strip()
    except ImportError:
        pass
    except Exception as e:
        print(f"  [PDF] PyPDF2 failed: {e}", flush=True)

    return ""


def _read_docx(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)
        return '\n'.join(parts).strip()
    except ImportError:
        print(f"  [DOCX] python-docx not installed. Run: pip install python-docx", flush=True)
        return ""
    except Exception as e:
        print(f"  [DOCX] Failed to read {file_path}: {e}", flush=True)
        return ""


def _read_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f"[Slide {slide_num}]\n" + '\n'.join(slide_texts))
        return '\n\n'.join(parts).strip()
    except ImportError:
        print(f"  [PPTX] python-pptx not installed. Run: pip install python-pptx", flush=True)
        return ""
    except Exception as e:
        print(f"  [PPTX] Failed to read {file_path}: {e}", flush=True)
        return ""


def _read_excel(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext in ('.xlsx', '.xlsm'):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_text = ' | '.join(str(c) for c in row if c is not None and str(c).strip())
                    if row_text:
                        rows.append(row_text)
                if rows:
                    parts.append(f"[Sheet: {sheet_name}]\n" + '\n'.join(rows))
            wb.close()
            return '\n\n'.join(parts).strip()
        except ImportError:
            print(f"  [EXCEL] openpyxl not installed.", flush=True)
        except Exception as e:
            print(f"  [EXCEL] openpyxl failed: {e}", flush=True)

    if ext == '.xls':
        try:
            import xlrd
            wb = xlrd.open_workbook(file_path)
            parts = []
            for sheet in wb.sheets():
                rows = []
                for row_idx in range(sheet.nrows):
                    row_text = ' | '.join(str(sheet.cell_value(row_idx, c))
                                         for c in range(sheet.ncols)
                                         if str(sheet.cell_value(row_idx, c)).strip())
                    if row_text:
                        rows.append(row_text)
                if rows:
                    parts.append(f"[Sheet: {sheet.name}]\n" + '\n'.join(rows))
            return '\n\n'.join(parts).strip()
        except ImportError:
            print(f"  [EXCEL] xlrd not installed.", flush=True)
        except Exception as e:
            print(f"  [EXCEL] xlrd failed: {e}", flush=True)

    if ext == '.csv':
        return _read_text_file(file_path)

    return ""


def _read_doc_legacy(file_path: str) -> str:
    try:
        import subprocess
        result = subprocess.run(['antiword', file_path], capture_output=True, text=True, timeout=10)
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
    except Exception:
        pass
    print(f"  [DOC] .doc format not fully supported.", flush=True)
    return ""


def _read_text_file(file_path: str) -> str:
    encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                content = f.read()
            if content and content.strip():
                return content.strip('\r\n').rstrip()
        except Exception:
            continue
    return ""


# ──────────────────────────────────────────────────────────────────────────────
# Main read function
# ──────────────────────────────────────────────────────────────────────────────
def read_file_safely(file_path: str, max_size: int = 500000) -> str:
    try:
        file_size = os.path.getsize(file_path)
        ext = Path(file_path).suffix.lower()

        if file_size == 0:
            return ""

        doc_max = max_size * 4 if ext in SUPPORTED_DOCUMENT_EXTENSIONS else max_size

        if file_size > doc_max:
            print(f"[SKIP] File too large ({file_size // 1024}KB): {file_path}")
            return ""

        if ext == '.pdf':
            content = _read_pdf(file_path)
        elif ext == '.docx':
            content = _read_docx(file_path)
        elif ext == '.doc':
            content = _read_doc_legacy(file_path)
        elif ext in ('.pptx', '.ppt'):
            content = _read_pptx(file_path)
        elif ext in ('.xlsx', '.xls', '.xlsm'):
            content = _read_excel(file_path)
        elif ext == '.csv':
            content = _read_text_file(file_path)
        else:
            content = _read_text_file(file_path)

        return content if content else ""

    except Exception as e:
        print(f"[SKIP] Could not read {file_path}: {str(e)}")
        return ""


def clean_text(text: str) -> str:
    text = text.replace('\x00', '')
    text = text.replace('\r\n', '\n')
    text = text.replace('\r', '\n')
    return text


def cleanup_directory(directory: str) -> None:
    try:
        if os.path.exists(directory):
            shutil.rmtree(directory)
    except Exception as e:
        print(f"Warning: Could not cleanup {directory}: {str(e)}")